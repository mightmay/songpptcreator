# -*- coding: utf-8 -*-
"""
Created on Sun May 27 00:00:04 2018

@author: Computer
"""
import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE

from flask import Flask, request, send_from_directory, send_file 
import xml.dom.minidom
from xml.dom.minidom import parse
from xml.dom import minidom
import xml.dom.minidom
from adjusttext import autoaddnewline
APP_ROOT = os.path.dirname(os.path.abspath(__file__)) 
lyricfile = os.path.join(APP_ROOT, 'songdata/amazing grace.xml')
savedirectory = os.path.join(APP_ROOT, 'finishedppt')
savefile = os.path.join(savedirectory, 'WorshipSongs.pptx')

songname= lyricfile

thai=1
english=0
mien=0
chord=1



def getsongdata(prs,songname,first_language,second_language,third_language,firsttextsizeint,secondtextsizeint,thirdtextsizeint,firsttextcolor,secondtextcolor,thirdtextcolor,chord):
    

    # create presentation
#    prs = Presentation()
    
    firsttextcolor = firsttextcolor
    if firsttextsizeint=='':
        
        firsttextsizeint=33
    firsttextsize=Pt(int(firsttextsizeint))
    
    secondtextcolor = secondtextcolor
    
    if secondtextsizeint=='':

        secondtextsizeint=33
    secondtextsize=Pt(int(secondtextsizeint))
    
    thirdtextcolor = thirdtextcolor
    
    if thirdtextsizeint=='':
        thirdtextsizeint=33
    thirdtextsize=Pt(int(thirdtextsizeint))
    
    # Open XML document containing song data using minidom parser
    songdata = minidom.parse(songname)

    
    # add song name to pptx slides
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    # get song ordering
    order = songdata.getElementsByTagName("order")[0]
    order_list=(order.childNodes[0].data)
    order_list = order_list.split(",")
    print(order_list)
    

    

    thirdlangnamestring="none"
    
    if third_language != "none":
        thirdlangxmlelement = third_language + "name"
        thirdlangname = songdata.getElementsByTagName(thirdlangxmlelement)[0]
        thirdlangnamestring=(thirdlangname.childNodes[0].data)

        if thirdlangnamestring == 'none':
            third_language= "none"
            
    if second_language != "none":
        secondlangxmlelement = second_language + "name"
        secondlangname = songdata.getElementsByTagName(secondlangxmlelement)[0]
        secondlangnamestring=(secondlangname.childNodes[0].data)
        subtitle.text = secondlangnamestring

        if secondlangnamestring == 'none':
            if thirdlangnamestring != 'none':
                second_language = third_language
                third_language= "none"
            if thirdlangnamestring == 'none':
                second_language="none"
                
            

    firstlangxmlelement = first_language + "name"
    firstlangname = songdata.getElementsByTagName(firstlangxmlelement)[0]
    firstlangnamestring=(firstlangname.childNodes[0].data)
    title.text = firstlangnamestring
    
    if firstlangnamestring == 'none':
            
        if secondlangnamestring != 'none' and thirdlangnamestring != 'none':

            first_language = second_language            
            second_language = third_language
            third_language ="none"
            
        if secondlangnamestring != 'none' and thirdlangnamestring == 'none':

            first_language = second_language
            second_language = "none"
        if secondlangnamestring == 'none' and thirdlangnamestring == 'none':
            return("no lyric")
            
            
    # delete "none" from Title and subtitle slide
    if subtitle.text == "none":
        subtitle.text = ""
    if title.text == "none":
        title.text = ""

    # count how many language
    language_count= 0
    if first_language != "none":
        language_count=language_count+1
    if second_language != "none":
        language_count=language_count+1
    if third_language != "none":
        language_count=language_count+1
                
    
    if language_count == 1:
        left1 = Inches(1)
        top1 = Inches(0.5)
        width1 = Inches(8)
        height1 = Inches(6)
               
            
    if language_count == 2:
        left1 = Inches(1)
        top1 = Inches(0.25)
        width1 = Inches(8)
        height1 = Inches(3.25)
        
        left2 = Inches(1)
        top2 = Inches(4)
        width2 = Inches(8)
        height2 = Inches(3.25)
        
    
    if language_count == 3:
        left1 = Inches(1)
        top1 = Inches(0.25)
        width1 = Inches(8)
        height1 = Inches(2)
        
        left2 = Inches(1)
        top2 = Inches(2.75)
        width2 = Inches(8)
        height2 = Inches(2)
        
        left3 = Inches(1)
        top3 = Inches(5)
        width3 = Inches(8)
        height3 = Inches(2)
    
    
    for versenum in order_list:
        vnum=str(versenum)
        print(vnum)
        

        if language_count==1:
            firstdata=songdata.getElementsByTagName(first_language)

            for first_itr in firstdata:
                firstlyric=first_itr.getElementsByTagName(vnum)[0]
                firstlyric=(firstlyric.childNodes[0].data)

                # create slide, add verse to slide
                blank_slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(blank_slide_layout)
                
                txBox1 = slide.shapes.add_textbox(left1, top1, width1, height1)
                tf1 = txBox1.text_frame
                para1 = tf1.add_paragraph()
                firstlyricaddedlinebreak=autoaddnewline(firstlyric,int(firsttextsizeint),language_count,first_language)
                print('font size: '+firsttextsizeint)
                print('line break: '+str(firstlyricaddedlinebreak))
                
                # no new line char added, then keep word wrap
                if(firstlyricaddedlinebreak==-1):
                    para1.text = firstlyric
                    tf1.word_wrap = True
                    
                # if new line char was added turn off word wrap
                else:
                    para1.text = firstlyricaddedlinebreak
                    tf1.word_wrap = False
                para1.font.size = firsttextsize
                
                para1.font.color.rgb = firsttextcolor
                para1.alignment=PP_ALIGN.CENTER
                

        if language_count==2:
            firstdata=songdata.getElementsByTagName(first_language)        
            seconddata=songdata.getElementsByTagName(second_language)
            for first_itr in firstdata:
                firstlyric=first_itr.getElementsByTagName(vnum)[0]
                firstlyric=(firstlyric.childNodes[0].data)
            for second_itr in seconddata:
                secondlyric=second_itr.getElementsByTagName(vnum)[0]
                secondlyric=(secondlyric.childNodes[0].data)

                # create slide, add verse to slide
                blank_slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(blank_slide_layout)
                
                txBox1 = slide.shapes.add_textbox(left1, top1, width1, height1)
                tf1 = txBox1.text_frame
                para1 = tf1.add_paragraph()

                para1.font.size = firsttextsize
                
                para1.font.color.rgb = firsttextcolor
                para1.alignment=PP_ALIGN.CENTER
                
                firstlyricaddedlinebreak=autoaddnewline(firstlyric,int(firsttextsizeint),language_count,first_language)
                print('font size: '+firsttextsizeint)
                print('line break: '+str(firstlyricaddedlinebreak))
                
                # no new line char added, then keep word wrap
                if(firstlyricaddedlinebreak==-1):
                    para1.text = firstlyric
                    tf1.word_wrap = True
                    
                # if new line char was added turn off word wrap
                else:
                    para1.text = firstlyricaddedlinebreak
                    tf1.word_wrap = False
                    
                txBox2 = slide.shapes.add_textbox(left2, top2, width2, height2)
                tf2 = txBox2.text_frame
                para2 = tf2.add_paragraph()

                para2.font.size = secondtextsize
                para2.font.color.rgb = secondtextcolor
                para2.alignment=PP_ALIGN.CENTER
                              
                secondlyricaddedlinebreak=autoaddnewline(secondlyric,int(secondtextsizeint),language_count,second_language)
                print('font size: '+secondtextsizeint)
                print('line break: '+str(secondlyricaddedlinebreak))
                
                # no new line char added, then keep word wrap
                if(secondlyricaddedlinebreak==-1):
                    para2.text = secondlyric
                    tf2.word_wrap = True
                    
                # if new line char was added turn off word wrap
                else:
                    para2.text = secondlyricaddedlinebreak
                    tf2.word_wrap = False
      
        if language_count==3:
            firstdata=songdata.getElementsByTagName(first_language)        
            seconddata=songdata.getElementsByTagName(second_language)        
            thirddata=songdata.getElementsByTagName(third_language)
            for first_itr in firstdata:
                firstlyric=first_itr.getElementsByTagName(vnum)[0]
                firstlyric=(firstlyric.childNodes[0].data)
            for second_itr in seconddata:
                secondlyric=second_itr.getElementsByTagName(vnum)[0]
                secondlyric=(secondlyric.childNodes[0].data)
            for third_itr in thirddata:
                thirdlyric=third_itr.getElementsByTagName(vnum)[0]
                thirdlyric=(thirdlyric.childNodes[0].data)
                
                # create slide, add verse to slide
                blank_slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(blank_slide_layout)
                
                txBox1 = slide.shapes.add_textbox(left1, top1, width1, height1)
                tf1 = txBox1.text_frame
                para1 = tf1.add_paragraph()
                para1.text = firstlyric
                tf1.word_wrap = True
                para1.font.size = firsttextsize
                para1.font.color.rgb = firsttextcolor
                para1.alignment=PP_ALIGN.CENTER
                
                txBox2 = slide.shapes.add_textbox(left2, top2, width2, height2)
                tf2 = txBox2.text_frame
                para2 = tf2.add_paragraph()
                para2.text = secondlyric
                tf2.word_wrap = True
                para2.font.size = secondtextsize
                para2.font.color.rgb = secondtextcolor
                para2.alignment=PP_ALIGN.CENTER
                
                txBox3 = slide.shapes.add_textbox(left3, top3, width3, height3)
                tf3 = txBox3.text_frame
                para3 = tf3.add_paragraph()
                para3.text = thirdlyric
                tf3.word_wrap = True
                para3.font.size = secondtextsize
                para3.font.color.rgb = thirdtextcolor
                para3.alignment=PP_ALIGN.CENTER       

                
#    prs.save(savefile)

                
def download_file():
    return send_from_directory(savedirectory,savefile, as_attachment=True)                

            

